import io
import email
from email import policy
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import extract_msg


# --- PDF ---
def _extract_pdf(uploaded_file) -> str:
    with pdfplumber.open(uploaded_file) as pdf:
        pages = [page.extract_text() or "" for page in pdf.pages]
    return "\n\n".join(pages).strip()


# --- Word ---
def _extract_docx(uploaded_file) -> str:
    doc = Document(uploaded_file)
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)
    return "\n\n".join(parts).strip()


# --- PowerPoint ---
def _extract_pptx(uploaded_file) -> str:
    pres = Presentation(uploaded_file)
    parts = []
    for i, slide in enumerate(pres.slides, start=1):
        slide_parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_parts.append(row_text)
        # Notities
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Notities: {notes}]")
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts).strip()


# --- Excel ---
def _extract_xlsx(uploaded_file) -> str:
    wb = load_workbook(uploaded_file, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts = [f"--- Werkblad: {sheet_name} ---"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                sheet_parts.append(" | ".join(cells))
        # Beperk extreem grote sheets
        if len(sheet_parts) > 500:
            sheet_parts = sheet_parts[:500] + ["[... werkblad ingekort ...]"]
        parts.append("\n".join(sheet_parts))
    return "\n\n".join(parts).strip()


# --- E-mail (.eml) ---
def _extract_eml(uploaded_file) -> str:
    if hasattr(uploaded_file, "read"):
        raw = uploaded_file.read()
        if isinstance(raw, str):
            raw = raw.encode("utf-8", errors="ignore")
    else:
        with open(uploaded_file, "rb") as f:
            raw = f.read()

    msg = email.message_from_bytes(raw, policy=policy.default)
    headers = [
        f"Van: {msg.get('From', '')}",
        f"Aan: {msg.get('To', '')}",
        f"Cc: {msg.get('Cc', '')}",
        f"Datum: {msg.get('Date', '')}",
        f"Onderwerp: {msg.get('Subject', '')}",
    ]
    headers = [h for h in headers if h.split(":", 1)[1].strip()]

    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                try:
                    body = part.get_content()
                    break
                except Exception:
                    pass
        if not body:
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    try:
                        body = part.get_content()
                        break
                    except Exception:
                        pass
    else:
        try:
            body = msg.get_content()
        except Exception:
            body = ""

    return "\n".join(headers) + "\n\n" + (body or "").strip()


# --- Outlook (.msg) ---
def _extract_msg(uploaded_file) -> str:
    if hasattr(uploaded_file, "read"):
        # extract_msg verwacht bytes of pad — schrijf tijdelijk naar bytes-buffer
        data = uploaded_file.read()
        buf = io.BytesIO(data)
        msg = extract_msg.Message(buf)
    else:
        msg = extract_msg.Message(uploaded_file)

    headers = [
        f"Van: {msg.sender or ''}",
        f"Aan: {msg.to or ''}",
        f"Cc: {msg.cc or ''}",
        f"Datum: {msg.date or ''}",
        f"Onderwerp: {msg.subject or ''}",
    ]
    headers = [h for h in headers if h.split(":", 1)[1].strip()]

    body = (msg.body or "").strip()
    return "\n".join(headers) + "\n\n" + body


# --- Dispatcher ---
_HANDLERS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".eml": _extract_eml,
    ".msg": _extract_msg,
}


def extract_text(uploaded_file) -> str:
    """Extract text from an uploaded Streamlit file object (PDF, Word, PPT, Excel, e-mail)."""
    name = getattr(uploaded_file, "name", "")
    suffix = Path(name).suffix.lower()

    handler = _HANDLERS.get(suffix)
    if handler:
        return handler(uploaded_file)

    # Fallback: probeer als PDF
    try:
        return _extract_pdf(uploaded_file)
    except Exception:
        uploaded_file.seek(0)
        return _extract_docx(uploaded_file)


def extract_text_from_path(path: str) -> str:
    """Extract text from a file on disk."""
    suffix = Path(path).suffix.lower()
    handler = _HANDLERS.get(suffix)
    if handler:
        return handler(path)
    return _extract_pdf(path)
